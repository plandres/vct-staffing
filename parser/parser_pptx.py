"""
PPTX parser for Seven2 VCT staffing slides.
Extracts staffing assignments from matrix tables in PPTX files.
"""

from pptx import Presentation
from pptx.util import Pt
from color_map import hex_to_workload, hex_to_status, normalize_initials, INITIALS_MAP
import re


def get_cell_fill_hex(cell) -> str | None:
    """Extract the fill color hex from a table cell."""
    try:
        fill = cell.fill
        if fill is None or fill.type is None:
            return None
        fc = fill.fore_color
        if fc and fc.type is not None:
            return str(fc.rgb)
    except Exception:
        pass
    return None


def get_text_color(cell) -> str | None:
    """Extract the dominant text color from a cell."""
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            try:
                if run.font.color and run.font.color.type is not None:
                    return str(run.font.color.rgb)
            except Exception:
                pass
    return None


def get_cell_text(cell) -> str:
    """Extract clean text from a cell."""
    return cell.text.strip()


def is_checkmark(text: str) -> bool:
    """Check if cell text is a checkmark (completed in previous periods)."""
    return text in {"u", "\u00fc", "\u2713", "\u2714"}


def parse_staffing_pptx(file_path: str) -> dict:
    """
    Parse a Seven2 VCT staffing PPTX and extract assignments.

    Returns:
        {
            "assignments": [
                {
                    "company": "Zwart",
                    "member_initials": "PLA",
                    "programs": ["Pricing", "Sales & M"],
                    "workload": "heavy",
                    "status": "ongoing",
                    "raw_text": "...",
                    "raw_fill_hex": "0E632A"
                },
                ...
            ],
            "companies_found": ["Zwart", "HRK Lunis", ...],
            "members_found": ["FC", "PLA", ...],
            "warnings": ["Unknown initials: XY", ...],
            "slide_count": 42
        }
    """
    prs = Presentation(file_path)
    assignments = []
    companies_found = set()
    members_found = set()
    warnings = []

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_table:
                continue

            table = shape.table
            rows = list(table.rows)
            cols = list(table.columns)

            if len(rows) < 2 or len(cols) < 2:
                continue

            # Try to detect if this is a staffing matrix
            # Strategy 1: Header row has member initials
            header_cells = [table.cell(0, c) for c in range(len(cols))]
            header_texts = [get_cell_text(cell) for cell in header_cells]

            # Check if headers contain known initials
            initials_in_header = []
            col_initials_map = {}  # col_index → initials

            for col_idx, text in enumerate(header_texts):
                if col_idx == 0:
                    continue  # Skip first column (usually company/program names)
                cleaned = normalize_initials(text)
                if cleaned in INITIALS_MAP:
                    initials_in_header.append(cleaned)
                    col_initials_map[col_idx] = cleaned

            # Strategy 2: First column has company names (detect by checking rows)
            row_labels = []
            for row_idx in range(1, len(rows)):
                cell = table.cell(row_idx, 0)
                text = get_cell_text(cell)
                if text:
                    row_labels.append((row_idx, text))

            # If we found initials in headers, this is likely a staffing matrix
            if len(initials_in_header) >= 2:
                for row_idx, company_name in row_labels:
                    companies_found.add(company_name)

                    for col_idx, initials in col_initials_map.items():
                        cell = table.cell(row_idx, col_idx)
                        text = get_cell_text(cell)
                        fill_hex = get_cell_fill_hex(cell)
                        text_color = get_text_color(cell)

                        workload = hex_to_workload(fill_hex)

                        # Skip empty cells, checkmarks, and "No VCT" cells
                        if workload == "none" or workload == "unknown":
                            continue
                        if is_checkmark(text):
                            continue
                        if not text and workload == "unknown":
                            continue

                        members_found.add(initials)

                        # Parse programs from text (comma-separated)
                        programs_list = []
                        if text:
                            # Split on comma, newline, or semicolon
                            parts = re.split(r"[,;\n]+", text)
                            programs_list = [
                                p.strip()
                                for p in parts
                                if p.strip() and not is_checkmark(p.strip())
                            ]

                        status = hex_to_status(text_color)

                        assignments.append(
                            {
                                "company": company_name,
                                "member_initials": initials,
                                "programs": programs_list if programs_list else [text or "General"],
                                "workload": workload,
                                "status": status,
                                "raw_text": text,
                                "raw_fill_hex": fill_hex,
                                "slide_index": slide_idx + 1,
                            }
                        )

    # Check for unknown initials in assignments
    for initials in members_found:
        if initials not in INITIALS_MAP:
            warnings.append(f"Unknown initials: {initials}")

    return {
        "assignments": assignments,
        "companies_found": sorted(companies_found),
        "members_found": sorted(members_found),
        "warnings": warnings,
        "slide_count": len(prs.slides),
    }
